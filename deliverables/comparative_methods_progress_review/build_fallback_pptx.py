#!/usr/bin/env python3
"""Build a stable PPTX from the visually approved PPTD page renders.

The public browser exporter can occasionally fail to emit a download.  This
fallback keeps the PPTD project as the editable source of truth and packages
the approved page renders into a 16:9 PowerPoint.  It starts from the project's
existing 16:9 embedded-font template so the resulting OpenXML package retains
font parts, then patches one fade transition into every slide.
"""

from __future__ import annotations

import importlib.util
import json
import re
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
REPORT_DIR = ROOT / "deliverables" / "comparative_methods_progress_review"
DECK_DIR = REPORT_DIR / "comparative_methods_progress_review_deck"
PAGES_DIR = DECK_DIR / ".qa-images" / "pages"
OUTPUT = DECK_DIR / "comparative_methods_progress_review.pptx"
TEMPLATE = ROOT / "docs" / "report" / "midterm_assessment_powerpoint_version" / "PPT模板.pptx"
EXPORTER = Path("/home/uais0/.codex/skills/open-kimi-ppt/scripts/export_pptx.py")


def load_exporter():
    spec = importlib.util.spec_from_file_location("open_kimi_export_pptx", EXPORTER)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Cannot load exporter helper: {EXPORTER}")
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst  # python-pptx has no public delete API.
    for slide_id in list(slide_ids):
        prs.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def page_number(path: Path) -> int:
    match = re.fullmatch(r"(\d+)\.jpe?g", path.name, flags=re.IGNORECASE)
    if not match:
        raise ValueError(f"Unexpected page filename: {path.name}")
    return int(match.group(1))


def package_summary(path: Path) -> dict[str, object]:
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        slides = sorted(
            name
            for name in names
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        )
        fade_count = 0
        transition_count = 0
        for name in slides:
            root = ET.fromstring(archive.read(name))
            transitions = root.findall(f"{{{p_ns}}}transition")
            transition_count += len(transitions)
            fade_count += sum(
                1 for node in transitions if node.find(f"{{{p_ns}}}fade") is not None
            )
        fonts = [
            name
            for name in names
            if name.startswith("ppt/fonts/") and not name.endswith("/")
        ]
        presentation = archive.read("ppt/presentation.xml")
        return {
            "file": str(path),
            "bytes": path.stat().st_size,
            "slides": len(slides),
            "transition_elements": transition_count,
            "fade_transitions": fade_count,
            "font_parts": len(fonts),
            "font_bytes": sum(archive.getinfo(name).file_size for name in fonts),
            "embedded_font_list": b"embeddedFontLst" in presentation,
            "zip_test": archive.testzip(),
        }


def main() -> int:
    pages = sorted(PAGES_DIR.glob("*.jpeg"), key=page_number)
    if [page_number(path) for path in pages] != list(range(1, 15)):
        raise RuntimeError(f"Expected approved page renders 1..14, got: {pages}")
    if not TEMPLATE.is_file():
        raise RuntimeError(f"Embedded-font template is missing: {TEMPLATE}")

    prs = Presentation(TEMPLATE)
    remove_all_slides(prs)
    # This embedded-font template has six custom layouts and no stock "Blank"
    # layout.  Its first layout is safe because the approved render is inserted
    # last and covers the full canvas.
    blank_layout = prs.slide_layouts[0]
    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(page),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    prs.save(OUTPUT)

    exporter = load_exporter()
    patched = exporter.patch_transitions(OUTPUT, "fade")
    summary = package_summary(OUTPUT)
    summary["transition_patched_slides"] = patched
    summary["source"] = "visually_approved_pptd_page_renders"

    expected = {
        "slides": 14,
        "transition_elements": 14,
        "fade_transitions": 14,
        "embedded_font_list": True,
        "zip_test": None,
    }
    for key, value in expected.items():
        if summary[key] != value:
            raise RuntimeError(f"Validation failed: {key}={summary[key]!r}, expected {value!r}")
    if int(summary["font_parts"]) < 1:
        raise RuntimeError("Validation failed: the PPTX has no embedded font parts")

    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
